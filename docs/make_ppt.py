"""
Онлайн дуудлага худалдааны цогц системийн хөгжүүлэлт
Дипломын ажлын танилцуулга — Т.Бөхбилэгт
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette (matches ExamGuard dark academic style) ────────────
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # very dark navy
ACCENT     = RGBColor(0x4F, 0x46, 0xE5)   # indigo-600
ACCENT2    = RGBColor(0x06, 0xB6, 0xD4)   # cyan-500
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def dark_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=DARK_BG)

def slide_number(slide, n, total=16):
    add_text(slide, f"{n} / {total}",
             W - Inches(1.2), H - Inches(0.4), Inches(1.1), Inches(0.35),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def section_tag(slide, text):
    add_rect(slide, Inches(0.5), Inches(0.3), Inches(2.4), Inches(0.38),
             fill=ACCENT)
    add_text(slide, text,
             Inches(0.55), Inches(0.3), Inches(2.3), Inches(0.38),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def divider(slide, y=Inches(1.05)):
    add_rect(slide, Inches(0.5), y, W - Inches(1.0), Pt(2), fill=ACCENT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)

# Left accent bar
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)

# Logo placeholder top-right
add_rect(s, W - Inches(2.5), Inches(0.3), Inches(2.0), Inches(0.9),
         fill=RGBColor(0x1E, 0x29, 0x4C))
add_text(s, "ШУТИС · МХТС", W - Inches(2.5), Inches(0.3),
         Inches(2.0), Inches(0.9), size=11, color=LIGHT_GRAY,
         align=PP_ALIGN.CENTER)

# Main title
add_text(s,
    "Онлайн дуудлага худалдааны\nцогц системийн хөгжүүлэлт",
    Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.2),
    size=38, bold=True, color=WHITE)

# Subtitle tag
add_rect(s, Inches(0.5), Inches(3.9), Inches(3.6), Inches(0.42), fill=ACCENT)
add_text(s, "Бакалаврын дипломын ажил",
         Inches(0.55), Inches(3.88), Inches(3.5), Inches(0.44),
         size=13, bold=True, color=WHITE)

# Info block
info = [
    ("Илтгэгч:",       "Т.Бөхбилэгт  (Түвшинзаяа Бөхбилэгт)"),
    ("Удирдагч:",      "Доктор (Ph.D), П.Энхтайван"),
    ("Мэргэжил:",      "Мэдээллийн технологи  ·  D061304"),
    ("Тэнхим:",        "Мэдээллийн технологийн тэнхим"),
    ("Хамгаалах он:",  "Улаанбаатар хот, 2026 он"),
]
y = Inches(4.5)
for label, val in info:
    add_text(s, label, Inches(0.6), y, Inches(2.1), Inches(0.36),
             size=13, bold=True, color=ACCENT2)
    add_text(s, val,   Inches(2.75), y, Inches(7.0), Inches(0.36),
             size=13, color=WHITE)
    y += Inches(0.42)

# Bottom line
add_rect(s, Inches(0.5), H - Inches(0.25), W - Inches(1.0),
         Pt(1.5), fill=ACCENT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
add_text(s, "АГУУЛГА", Inches(0.5), Inches(0.28), W, Inches(0.55),
         size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 2)

items = [
    ("01", "Судалгааны зорилго ба ач холбогдол"),
    ("02", "Онолын үндэслэл ба харьцуулсан судалгаа"),
    ("03", "Системийн архитектур ба технологи"),
    ("04", "Хэрэгжүүлэлт — Backend, Веб, Мобайл"),
    ("05", "Өгөгдлийн загвар (ERD)"),
    ("06", "Туршилт ба үр дүн"),
    ("07", "Дүгнэлт ба цаашдын чиглэл"),
]

cols = [
    (Inches(0.5),  items[:4]),
    (Inches(6.8),  items[4:]),
]
for cx, grp in cols:
    y = Inches(1.3)
    for num, title in grp:
        add_rect(s, cx, y, Inches(0.62), Inches(0.56),
                 fill=ACCENT)
        add_text(s, num, cx, y, Inches(0.62), Inches(0.56),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, cx + Inches(0.72), y,
                 Inches(5.6), Inches(0.56), size=15, color=WHITE)
        y += Inches(0.82)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH OBJECTIVE
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "01 · ЗОРИЛГО")
add_text(s, "СУДАЛГААНЫ ЗОРИЛГО", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 3)

# Problem box
add_rect(s, Inches(0.5), Inches(1.2), Inches(5.9), Inches(2.5),
         fill=RGBColor(0x1E, 0x29, 0x4C), line=ACCENT)
add_text(s, "Асуудал", Inches(0.65), Inches(1.25),
         Inches(5.6), Inches(0.4), size=14, bold=True, color=ACCENT2)
problems = [
    "• Монголд бодит цагийн дуудлага худалдааны систем байхгүй",
    "• Одоогийн платформ (Unegui, Shoppy) зөвхөн тогтмол үнийн загвартай",
    "• Мобайл хэрэглэгчдэд зориулсан cross-platform шийдэл дутагдалтай",
    "• Үнэ тогтоолт зах зээлийн зарчмаар явагдах боломжгүй",
]
y = Inches(1.72)
for p in problems:
    add_text(s, p, Inches(0.65), y, Inches(5.6), Inches(0.38),
             size=12, color=LIGHT_GRAY)
    y += Inches(0.37)

# Goal box
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.5),
         fill=RGBColor(0x1E, 0x29, 0x4C), line=GREEN)
add_text(s, "Зорилго", Inches(6.95), Inches(1.25),
         Inches(5.7), Inches(0.4), size=14, bold=True, color=GREEN)
goals = [
    "✓  Бодит цагийн (Socket.IO) дуудлага худалдааны систем",
    "✓  Web + iOS + Android — нэг кодоор 3 платформ",
    "✓  QPay Монгол төлбөрийн систем интеграц",
    "✓  AI ангилал санал, итгэлцлийн оноо (Trust Score)",
]
y = Inches(1.72)
for g in goals:
    add_text(s, g, Inches(6.95), y, Inches(5.7), Inches(0.38),
             size=12, color=WHITE)
    y += Inches(0.37)

# Research question
add_rect(s, Inches(0.5), Inches(3.9), W - Inches(1.0), Inches(0.72),
         fill=ACCENT)
add_text(s,
    "Судалгааны үндсэн асуудал:  Монголын зах зээлд тохирсон, бодит цагийн,"
    " олон платформ дээр ажилладаг онлайн дуудлага худалдааны системийг хэрхэн хөгжүүлэх вэ?",
    Inches(0.7), Inches(3.92), W - Inches(1.4), Inches(0.68),
    size=13, bold=True, color=WHITE)

# Stats row
stats = [
    ("3", "Платформ\n(Web/iOS/Android)"),
    ("34+", "REST API\nEndpoint"),
    ("11", "MongoDB\nCollection"),
    ("66", "Ангилал"),
    ("12K+", "Кодын мөр"),
]
x = Inches(0.5)
for val, lbl in stats:
    add_rect(s, x, Inches(4.8), Inches(2.3), Inches(1.4),
             fill=RGBColor(0x1E, 0x29, 0x4C))
    add_text(s, val, x, Inches(4.85), Inches(2.3), Inches(0.65),
             size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x, Inches(5.52), Inches(2.3), Inches(0.6),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    x += Inches(2.55)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — THEORY & COMPARISON
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "02 · ОНОЛ")
add_text(s, "ХАРЬЦУУЛСАН СУДАЛГАА", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 4)

# Auction types left
add_text(s, "Дуудлага худалдааны төрлүүд",
         Inches(0.5), Inches(1.15), Inches(4.0), Inches(0.4),
         size=14, bold=True, color=ACCENT2)
types = [
    ("🔨 Англи (English)", "Үнэ доороос өсдөг · Манай систем"),
    ("📉 Голланд (Dutch)", "Үнэ дээрээс буурдаг · BuyNow-ээр"),
    ("🔒 Битэр (Sealed-bid)", "Нууц санал · Хугацаа дуусахад нээгдэнэ"),
    ("💡 Виккери (Vickrey)", "2-р өндөр үнэ · Google AdWords загвар"),
]
y = Inches(1.65)
for title, desc in types:
    add_rect(s, Inches(0.5), y, Inches(5.6), Inches(0.68),
             fill=RGBColor(0x1E, 0x29, 0x4C))
    add_text(s, title, Inches(0.65), y + Inches(0.04),
             Inches(5.3), Inches(0.32), size=13, bold=True, color=WHITE)
    add_text(s, desc, Inches(0.65), y + Inches(0.34),
             Inches(5.3), Inches(0.28), size=11, color=LIGHT_GRAY)
    y += Inches(0.76)

# Comparison table right
headers = ["Шинж чанар", "Yahoo!", "Mercari", "eBay", "Манай"]
rows = [
    ["Auction type",  "English", "Fixed",  "Multi",  "Eng+Buy"],
    ["Real-time bid", "✓",       "—",      "✓",      "✓"],
    ["Mobile parity", "✓",       "✓✓",     "✓",      "✓✓"],
    ["AI ангилал",    "✓",       "✓",      "✓",      "✓"],
    ["Дотоод төлбөр", "Y.Wallet","M.Pay",  "PayPal", "QPay ✓"],
    ["Монгол хэл",    "—",       "—",      "—",      "✓"],
]
col_w = [Inches(1.9), Inches(0.9), Inches(0.9), Inches(0.7), Inches(1.0)]
xs = [Inches(6.5)]
for cw in col_w[:-1]:
    xs.append(xs[-1] + cw)

# Header row
for i, (h, cw, x) in enumerate(zip(headers, col_w, xs)):
    bg = ACCENT if i == len(headers)-1 else RGBColor(0x1E, 0x29, 0x4C)
    add_rect(s, x, Inches(1.15), cw, Inches(0.42), fill=bg)
    add_text(s, h, x + Inches(0.05), Inches(1.17),
             cw - Inches(0.1), Inches(0.38),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    ry = Inches(1.57) + ri * Inches(0.62)
    for ci, (cell, cw, x) in enumerate(zip(row, col_w, xs)):
        bg = RGBColor(0x16, 0x21, 0x38) if ri % 2 == 0 else RGBColor(0x1E, 0x29, 0x4C)
        if ci == len(row) - 1:
            bg = RGBColor(0x1A, 0x1F, 0x4C)
        add_rect(s, x, ry, cw, Inches(0.58), fill=bg)
        col = GREEN if (ci == len(row)-1 and cell not in ["Eng+Buy","✓✓","QPay ✓","✓"]) else (GREEN if cell in ["✓","✓✓"] else (RED if cell == "—" else WHITE))
        if ci == len(row) - 1:
            col = GOLD
        add_text(s, cell, x + Inches(0.04), ry + Inches(0.12),
                 cw - Inches(0.08), Inches(0.36),
                 size=11, color=col, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "03 · АРХИТЕКТУР")
add_text(s, "СИСТЕМИЙН АРХИТЕКТУР", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 5)

layers = [
    ("CLIENT LAYER",
     "React.js + Vite  ·  React Native + Expo SDK 54  ·  Expo Router v6",
     RGBColor(0x1E, 0x29, 0x4C), ACCENT2),
    ("APPLICATION LAYER",
     "Node.js + Express.js  ·  Socket.IO  ·  JWT  ·  Cron Jobs  ·  Firebase Admin",
     RGBColor(0x14, 0x24, 0x1A), GREEN),
    ("DATA LAYER",
     "MongoDB Atlas (11 collection)  ·  Cloudinary CDN  ·  QPay  ·  Firebase FCM",
     RGBColor(0x24, 0x1A, 0x10), GOLD),
]

y = Inches(1.2)
for title, desc, bg, col in layers:
    add_rect(s, Inches(0.5), y, Inches(8.5), Inches(1.35), fill=bg, line=col)
    add_text(s, title, Inches(0.7), y + Inches(0.12),
             Inches(8.0), Inches(0.42), size=16, bold=True, color=col)
    add_text(s, desc, Inches(0.7), y + Inches(0.58),
             Inches(8.0), Inches(0.56), size=13, color=WHITE)
    if y < Inches(4.5):
        add_text(s, "⬇  REST / WebSocket" if y < Inches(3.0) else "⬇  Mongoose / SDK",
                 Inches(3.5), y + Inches(1.38), Inches(3.0), Inches(0.36),
                 size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    y += Inches(1.75)

# Tech tags right side
tech = ["Node.js v20", "Express.js", "MongoDB Atlas", "Socket.IO v4",
        "React.js", "Vite", "Tailwind CSS", "React Native",
        "Expo SDK 54", "Firebase", "Cloudinary", "QPay"]
tx = Inches(9.2)
ty = Inches(1.2)
for t in tech:
    add_rect(s, tx, ty, Inches(3.8), Inches(0.38),
             fill=RGBColor(0x1E, 0x29, 0x4C))
    add_text(s, "▸  " + t, tx + Inches(0.1), ty + Inches(0.05),
             Inches(3.6), Inches(0.28), size=12, color=WHITE)
    ty += Inches(0.46)
    if ty > Inches(6.5):
        break

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — BACKEND IMPLEMENTATION
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "04 · ХЭРЭГЖҮҮЛЭЛТ")
add_text(s, "BACKEND ХӨГЖҮҮЛЭЛТ", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 6)

modules = [
    ("User Module",    "Бүртгэл, нэвтрэлт, Google OAuth, Firebase Phone Auth, JWT"),
    ("Product Module", "CRUD, зургийн upload (Cloudinary), AI ангилал санал"),
    ("Bidding Module", "Atomic transaction, race condition урьдчилан сэргийлэх"),
    ("Cron Job",       "1 минут тутамд дуудлага эхлүүлэх / дуусгах, ялагч тогтоох"),
    ("Socket.IO",      "bidUpdate, outbid, auctionEnded, notification event-үүд"),
    ("QPay",           "Invoice үүсгэх, QR харуулах, webhook-оор баланс шинэчлэх"),
]
y = Inches(1.15)
for i, (mod, desc) in enumerate(modules):
    x = Inches(0.5) if i % 2 == 0 else Inches(6.8)
    if i % 2 == 0 and i > 0:
        y += Inches(1.05)
    add_rect(s, x, y, Inches(6.0), Inches(0.92),
             fill=RGBColor(0x1E, 0x29, 0x4C), line=ACCENT)
    add_text(s, mod, x + Inches(0.15), y + Inches(0.08),
             Inches(5.7), Inches(0.36), size=14, bold=True, color=ACCENT2)
    add_text(s, desc, x + Inches(0.15), y + Inches(0.48),
             Inches(5.7), Inches(0.36), size=12, color=LIGHT_GRAY)

# Socket events table
add_text(s, "Үндсэн Socket.IO event-үүд",
         Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.4),
         size=13, bold=True, color=ACCENT2)
events = [
    ("joinAuction",  "Client→Server", "Дуудлагын room-д орох"),
    ("placeBid",     "Client→Server", "Шинэ үнэ санал"),
    ("bidUpdate",    "Server→Room",   "Бүх оролцогчид шинэ үнэ зарлах"),
    ("auctionEnded", "Server→Room",   "Дуудлага дуусгасан, ялагч зарлах"),
    ("outbid",       "Server→User",   "Гүйцсэн мэдэгдэл → push notification"),
]
ew = [Inches(2.2), Inches(1.8), Inches(7.8)]
ex = [Inches(0.5), Inches(2.75), Inches(4.6)]
for ri, row in enumerate(events):
    ry = Inches(4.82) + ri * Inches(0.46)
    for ci, (cell, cw, x) in enumerate(zip(row, ew, ex)):
        bg = RGBColor(0x16, 0x21, 0x38) if ri % 2 == 0 else RGBColor(0x1A, 0x26, 0x42)
        add_rect(s, x, ry, cw, Inches(0.42), fill=bg)
        col = ACCENT2 if ci == 0 else (GOLD if ci == 1 else WHITE)
        add_text(s, cell, x + Inches(0.06), ry + Inches(0.08),
                 cw - Inches(0.1), Inches(0.28), size=11, color=col)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — WEB & MOBILE UI
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "04 · ХЭРЭГЖҮҮЛЭЛТ")
add_text(s, "ВЕБ БА МОБАЙЛ АППЛИКЕЙШН", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 7)

# Web column
add_rect(s, Inches(0.5), Inches(1.15), Inches(5.9), Inches(0.42), fill=ACCENT)
add_text(s, "Веб аппликейшн (React.js + Vite)",
         Inches(0.6), Inches(1.17), Inches(5.7), Inches(0.38),
         size=13, bold=True, color=WHITE)
web_feats = [
    "🏠  Нүүр хуудас — категориор шүүх, хайлт, auction feed",
    "📦  Бараа нэмэх — 20 зураг upload, AI ангилал санал, ноорог",
    "🔨  Бараа дэлгэрэнгүй — countdown таймер, bid history, шууд авах",
    "⚙️   Админ самбар — хэрэглэгч, ангилал, гүйлгээ, статистик",
    "👤  Профайл — баланс, QPay цэнэглэлт, зарлагуудын түүх",
    "🌓  Dark/Light горим, i18n (Монгол/Англи)",
]
y = Inches(1.65)
for f in web_feats:
    add_text(s, f, Inches(0.65), y, Inches(5.6), Inches(0.52),
             size=12, color=WHITE)
    y += Inches(0.54)

# Mobile column
add_rect(s, Inches(6.9), Inches(1.15), Inches(6.0), Inches(0.42), fill=GREEN)
add_text(s, "Мобайл аппликейшн (React Native + Expo)",
         Inches(7.0), Inches(1.17), Inches(5.8), Inches(0.38),
         size=13, bold=True, color=DARK_BG)
mobile_feats = [
    "📱  iOS + Android — нэг кодоор 2 платформ",
    "🔐  Google OAuth + Firebase Phone Auth (SMS OTP)",
    "📋  Регистрийн дугаар — Кирилл үсгийн picker",
    "📸  expo-image-picker — зураг авах/сонгох",
    "⏱   Socket.IO real-time bid update, countdown",
    "💳  QPay QR код — 3 секунд тутам polling",
    "🔔  FCM push notification — outbid, won мэдэгдэл",
    "⭐  Trust Score, Watchlist, Review системүүд",
]
y = Inches(1.65)
for f in mobile_feats:
    add_text(s, f, Inches(7.05), y, Inches(5.7), Inches(0.46),
             size=12, color=WHITE)
    y += Inches(0.48)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — ERD
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "05 · ӨГӨГДЛИЙН ЗАГВАР")
add_text(s, "ӨГӨГДЛИЙН САНГИЙН ЗАГВАР (ERD)", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 8)

collections = [
    ("User",         "name, email, password(bcrypt),\nrole, balance, trustScore, fcmTokens",  ACCENT),
    ("Product",      "user, category, price, currentBid,\nauctionStatus, bidDeadline, sellType", GREEN),
    ("Category",     "parent(self-ref), title, slug,\nfieldSchema[]",                          GOLD),
    ("Bidding",      "user, product, price, createdAt",                                        RGBColor(0x88,0x4E,0xA0)),
    ("Watchlist",    "user, product, notifyOnStart",                                           RED),
    ("Transaction",  "buyer, seller, product, amount",                                         RGBColor(0x1E,0x84,0x49)),
    ("Review",       "product, fromUser, toUser, rating",                                      RGBColor(0xE6,0x7E,0x22)),
    ("Notification", "user, product, type, read",                                              RGBColor(0x1A,0x52,0x76)),
]
positions = [
    (Inches(0.5),  Inches(1.2)),
    (Inches(0.5),  Inches(2.85)),
    (Inches(3.5),  Inches(1.2)),
    (Inches(3.5),  Inches(2.85)),
    (Inches(3.5),  Inches(4.5)),
    (Inches(6.5),  Inches(1.2)),
    (Inches(6.5),  Inches(2.85)),
    (Inches(9.5),  Inches(1.2)),
]
for (name, fields, col), (px, py) in zip(collections, positions):
    add_rect(s, px, py, Inches(2.75), Inches(1.35),
             fill=RGBColor(0x1E, 0x29, 0x4C), line=col)
    add_rect(s, px, py, Inches(2.75), Inches(0.38), fill=col)
    add_text(s, name, px + Inches(0.08), py + Inches(0.04),
             Inches(2.6), Inches(0.32), size=13, bold=True,
             color=DARK_BG if col == GOLD else WHITE)
    add_text(s, fields, px + Inches(0.08), py + Inches(0.42),
             Inches(2.6), Inches(0.88), size=10, color=LIGHT_GRAY)

# Relationships note
add_rect(s, Inches(9.5), Inches(2.85), Inches(3.5), Inches(3.8),
         fill=RGBColor(0x14, 0x1E, 0x38))
add_text(s, "Хамаарлууд", Inches(9.65), Inches(2.9),
         Inches(3.2), Inches(0.38), size=13, bold=True, color=ACCENT2)
rels = [
    "User  →  Product (1:N, seller)",
    "Category  →  Product (1:N)",
    "Category  →  Category (self, parent)",
    "User  →  Bidding (1:N)",
    "Product  →  Bidding (1:N)",
    "User  →  Watchlist (1:N)",
    "User  →  Transaction (1:N)",
    "Product  →  Transaction (1:1)",
    "User  →  Review (1:N)",
    "User  →  Notification (1:N)",
]
ry = Inches(3.35)
for r in rels:
    add_text(s, "▸ " + r, Inches(9.65), ry,
             Inches(3.2), Inches(0.35), size=10, color=WHITE)
    ry += Inches(0.36)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — TESTING
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "06 · ТУРШИЛТ")
add_text(s, "ТУРШИЛТ БА ҮР ДҮН", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 9)

# API test summary
add_text(s, "API Туршилт (Insomnia) — 34 endpoint, бүгд тэнцсэн",
         Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.42),
         size=14, bold=True, color=ACCENT2)
api_rows = [
    ("POST /api/users/register",   "201", "Бүртгэл амжилттай"),
    ("POST /api/users/login",      "200", "JWT токен буцсан"),
    ("POST /api/bidding",          "201", "Bid бүртгэгдсэн"),
    ("POST /api/bidding (бага үнэ)","400","Bid too low — зөв алдаа"),
    ("POST /api/product/:id/buy-now","200","Шууд худалдаж авсан"),
    ("GET  /api/product/products", "200", "Бараан жагсаалт буцсан"),
]
cols_w = [Inches(4.0), Inches(0.8), Inches(7.6)]
cols_x = [Inches(0.5), Inches(4.55), Inches(5.4)]
for ri, row in enumerate(api_rows):
    ry = Inches(1.65) + ri * Inches(0.46)
    for ci, (cell, cw, cx) in enumerate(zip(row, cols_w, cols_x)):
        bg = RGBColor(0x16,0x21,0x38) if ri%2==0 else RGBColor(0x1A,0x26,0x42)
        add_rect(s, cx, ry, cw, Inches(0.42), fill=bg)
        col = ACCENT2 if ci==0 else (GREEN if cell=="200" or cell=="201" else (RED if cell=="400" else WHITE))
        add_text(s, cell, cx+Inches(0.06), ry+Inches(0.08),
                 cw-Inches(0.1), Inches(0.28), size=11, color=col)

# Performance metrics
add_text(s, "Гүйцэтгэлийн үзүүлэлт",
         Inches(0.5), Inches(4.6), Inches(6.0), Inches(0.4),
         size=14, bold=True, color=ACCENT2)
perf = [
    ("GET /api/product",  "85 мс",  "220 мс"),
    ("POST /api/bidding", "110 мс", "290 мс"),
    ("POST login",        "180 мс", "380 мс"),
]
ph = ["Endpoint","Avg latency","P95 latency"]
px2 = [Inches(0.5), Inches(4.5), Inches(6.5)]
pw  = [Inches(3.9), Inches(1.9), Inches(1.9)]
for ci, (h, x, w) in enumerate(zip(ph, px2, pw)):
    add_rect(s, x, Inches(5.08), w, Inches(0.36), fill=ACCENT)
    add_text(s, h, x+Inches(0.05), Inches(5.1), w-Inches(0.1),
             Inches(0.3), size=11, bold=True, color=WHITE)
for ri, row in enumerate(perf):
    ry = Inches(5.48) + ri*Inches(0.4)
    for ci, (cell, x, w) in enumerate(zip(row, px2, pw)):
        bg = RGBColor(0x16,0x21,0x38) if ri%2==0 else RGBColor(0x1A,0x26,0x42)
        add_rect(s, x, ry, w, Inches(0.36), fill=bg)
        col = GREEN if ci > 0 else WHITE
        add_text(s, cell, x+Inches(0.05), ry+Inches(0.06),
                 w-Inches(0.08), Inches(0.26), size=11, color=col)

# Right — unit test coverage
add_text(s, "Тестийн бүрхэвч",
         Inches(8.5), Inches(4.6), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=ACCENT2)
cov_items = [
    ("Backend line coverage", "65%",  GOLD),
    ("Controller coverage",   "80%+", GREEN),
    ("Integration тест",      "18",   ACCENT2),
    ("Real-time bid latency",  "<1 сек", GREEN),
    ("P95 API latency",       "<400 мс", GREEN),
]
cy = Inches(5.1)
for lbl, val, col in cov_items:
    add_rect(s, Inches(8.5), cy, Inches(4.5), Inches(0.42),
             fill=RGBColor(0x1E,0x29,0x4C))
    add_text(s, lbl, Inches(8.65), cy+Inches(0.08),
             Inches(3.0), Inches(0.28), size=12, color=WHITE)
    add_text(s, val, Inches(11.5), cy+Inches(0.08),
             Inches(1.4), Inches(0.28), size=13, bold=True,
             color=col, align=PP_ALIGN.RIGHT)
    cy += Inches(0.5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "06 · АЮУЛГҮЙ БАЙДАЛ")
add_text(s, "АЮУЛГҮЙ БАЙДАЛ (OWASP Top-10)",
         Inches(0.5), Inches(0.28), W, Inches(0.55),
         size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 10)

owasp = [
    ("A01","Broken Access Control",   "Role-based middleware, ownership check"),
    ("A02","Cryptographic Failures",  "bcrypt 10 round, HTTPS only"),
    ("A03","Injection",               "Mongoose strict mode, input sanitization"),
    ("A05","Security Misconfiguration","helmet, .env, no debug in prod"),
    ("A06","Vulnerable Components",   "npm audit, Dependabot"),
    ("A07","Auth Failures",           "express-rate-limit, account lockout"),
    ("A09","Logging Failures",        "Winston logs, centralized error handler"),
]
ow = [Inches(0.65), Inches(2.2), Inches(4.2)]
ox = [Inches(0.5),  Inches(1.2), Inches(3.45)]
y = Inches(1.2)
for row in owasp:
    for cell, x, w in zip(row, ox, ow):
        bg = RGBColor(0x1A, 0x1F, 0x38) if row[0] else DARK_BG
        add_rect(s, x, y, w, Inches(0.62), fill=bg)
        col = GOLD if x == ox[0] else (ACCENT2 if x == ox[1] else GREEN)
        add_text(s, cell, x+Inches(0.08), y+Inches(0.14),
                 w-Inches(0.14), Inches(0.36), size=12,
                 bold=(x==ox[0]), color=col)
    y += Inches(0.7)

# Auth flow right
add_text(s, "Нэвтрэлтийн урсгал",
         Inches(8.5), Inches(1.15), Inches(4.6), Inches(0.4),
         size=14, bold=True, color=ACCENT2)
auth_steps = [
    ("1", "И-мэйл баталгаажуулах (OTP код)"),
    ("2", "EULA зөвшөөрөх"),
    ("3", "Бүртгэл дуусгах"),
    ("4", "JWT токен олгох (7 хоног)"),
    ("",  ""),
    ("✓", "Google OAuth 2.0"),
    ("✓", "Firebase Phone Auth (SMS)"),
    ("✓", "Refresh Token (AsyncStorage)"),
]
ay = Inches(1.65)
for num, step in auth_steps:
    if not num:
        ay += Inches(0.2)
        continue
    add_rect(s, Inches(8.5), ay, Inches(0.45), Inches(0.42),
             fill=ACCENT if num.isdigit() else GREEN)
    add_text(s, num, Inches(8.5), ay, Inches(0.45), Inches(0.42),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, Inches(9.05), ay+Inches(0.06),
             Inches(3.9), Inches(0.32), size=12, color=WHITE)
    ay += Inches(0.52)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — OBJECTIVES MET
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "06 · ҮР ДҮН")
add_text(s, "ЗОРИЛТЫН БИЕЛЭЛТ", Inches(0.5), Inches(0.28),
         W, Inches(0.55), size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 11)

obj_rows = [
    ("✓", "Бодит цагийн дуудлага",   "Socket.IO bidUpdate < 1 сек"),
    ("✓", "QPay интеграц",            "Sandbox тест, архитектур бэлэн"),
    ("✓", "Google OAuth",             "Веб + Мобайл хоёуланд хэрэгжсэн"),
    ("✓", "Автомат дуудлага",         "Cron job, ялагч тогтоох, balance"),
    ("✓", "MongoDB Atlas",            "11 collection, replicaset"),
    ("✓", "Аюулгүй байдал",           "JWT+bcrypt+helmet, P95<400мс"),
    ("✓", "Тест",                     "Unit+Integration+API+Mobile"),
]
y = Inches(1.2)
for icon, obj, detail in obj_rows:
    add_rect(s, Inches(0.5), y, Inches(0.52), Inches(0.68), fill=GREEN)
    add_text(s, icon, Inches(0.5), y, Inches(0.52), Inches(0.68),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.1), y, Inches(4.5), Inches(0.68),
             fill=RGBColor(0x1E,0x29,0x4C))
    add_text(s, obj, Inches(1.2), y+Inches(0.14),
             Inches(4.3), Inches(0.36), size=13, bold=True, color=WHITE)
    add_rect(s, Inches(5.7), y, Inches(7.3), Inches(0.68),
             fill=RGBColor(0x14,0x1E,0x38))
    add_text(s, detail, Inches(5.85), y+Inches(0.14),
             Inches(7.1), Inches(0.36), size=13, color=LIGHT_GRAY)
    y += Inches(0.8)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)
section_tag(s, "07 · ДҮГНЭЛТ")
add_text(s, "ДҮГНЭЛТ БА ЦААШДЫН ЧИГЛЭЛ",
         Inches(0.5), Inches(0.28), W, Inches(0.55),
         size=28, bold=True, color=WHITE)
divider(s)
slide_number(s, 12)

# Achievements
add_rect(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.42), fill=GREEN)
add_text(s, "✅  Гол ололтууд", Inches(0.6), Inches(1.17),
         Inches(5.8), Inches(0.38), size=14, bold=True, color=DARK_BG)
achievements = [
    "▸  Multi-platform систем: Web + iOS + Android нэгэн зэрэг",
    "▸  Бодит цагийн дуудлага: Socket.IO <1 сек update",
    "▸  Монгол зах зээлд тохирсон: QPay, 66 категори, Кирилл бүртгэл",
    "▸  AI технологи: Ангилал автоматаар санал болгох",
    "▸  Trust Score: Хэрэглэгчийн найдвартай байдлыг тоогоор хэмжих",
    "▸  Эдийн засгийн онол: Vickrey-Myerson, Akerlof-г практикт хэрэгжүүлэв",
]
y = Inches(1.65)
for a in achievements:
    add_text(s, a, Inches(0.65), y, Inches(5.7), Inches(0.48),
             size=12, color=WHITE)
    y += Inches(0.5)

# Future work
add_rect(s, Inches(7.0), Inches(1.15), Inches(6.0), Inches(0.42),
         fill=RGBColor(0x88,0x4E,0xA0))
add_text(s, "🔮  Цаашдын хөгжүүлэлт", Inches(7.1), Inches(1.17),
         Inches(5.8), Inches(0.38), size=14, bold=True, color=WHITE)
future = [
    "▸  Live video streaming (WebRTC / Twilio)",
    "▸  QPay бодит орчинд deploy, merchant agreement",
    "▸  AI chatbot (Claude/OpenAI) — хэрэглэгч тусламж",
    "▸  Anti-sniping — сүүлийн 60 сек +30 сек сунгах",
    "▸  Blockchain / NFT auction, smart contract",
    "▸  Multi-currency (USD, CNY), бодит цагийн ханш",
    "▸  Penetration testing (OWASP ASVS)",
]
y = Inches(1.65)
for f in future:
    add_text(s, f, Inches(7.1), y, Inches(5.8), Inches(0.48),
             size=12, color=WHITE)
    y += Inches(0.5)

# Bottom stat strip
stats2 = [
    ("3", "Платформ"), ("11", "Collection"), ("34+", "Endpoint"),
    ("7", "Socket event"), ("66", "Ангилал"), ("12K+", "Кодын мөр"),
    ("<400мс", "P95 latency"), ("<1сек", "Real-time"),
]
x = Inches(0.5)
bw = (W - Inches(1.0)) / len(stats2)
for val, lbl in stats2:
    add_rect(s, x, Inches(6.3), bw - Inches(0.05), Inches(0.9),
             fill=RGBColor(0x1E,0x29,0x4C))
    add_text(s, val, x, Inches(6.32), bw, Inches(0.45),
             size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x, Inches(6.76), bw, Inches(0.36),
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    x += bw

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
add_rect(s, 0, 0, Inches(0.18), H, fill=ACCENT)

# Big decorative circle
add_rect(s, Inches(7.5), Inches(-1.0), Inches(6.5), Inches(6.5),
         fill=RGBColor(0x1E,0x29,0x4C))

add_text(s, "АНХААРАЛТАЙ\nСОНССОНД БАЯРЛАЛАА",
         Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.0),
         size=42, bold=True, color=WHITE)

add_rect(s, Inches(0.5), Inches(3.7), Inches(4.0), Inches(0.05), fill=ACCENT)

info2 = [
    ("Гүйцэтгэгч:",  "Т.Бөхбилэгт  (Түвшинзаяа Бөхбилэгт)"),
    ("Удирдагч:",    "Доктор (Ph.D), П.Энхтайван"),
    ("Имэйл:",       "buhuu1125@gmail.com"),
    ("ШУТИС МХТС ·", "Мэдээллийн технологийн тэнхим · 2026"),
]
y = Inches(3.95)
for lbl, val in info2:
    add_text(s, lbl, Inches(0.5), y, Inches(2.0), Inches(0.4),
             size=13, bold=True, color=ACCENT2)
    add_text(s, val,  Inches(2.55), y, Inches(6.0), Inches(0.4),
             size=13, color=WHITE)
    y += Inches(0.46)

add_text(s,
    "Онлайн дуудлага худалдааны цогц системийн хөгжүүлэлт",
    Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.5),
    size=14, color=LIGHT_GRAY)

add_rect(s, Inches(0.5), H - Inches(0.25), W - Inches(1.0),
         Pt(1.5), fill=ACCENT)

# ── Save ──────────────────────────────────────────────────────────────
out = r"C:\Users\bukhbtu01\Downloads\onlineauction-clean\docs\thesis_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
